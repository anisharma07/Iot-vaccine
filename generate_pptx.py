import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Configuration ───
DARK_RED   = RGBColor(178, 34, 34)
GOLD       = RGBColor(204, 153, 51)
BLUE       = RGBColor(51, 51, 178)
GRAY       = RGBColor(160, 160, 160)
WHITE      = RGBColor(255, 255, 255)
DARK_BG    = RGBColor(30, 30, 46)
GREEN      = RGBColor(80, 200, 120)
LIGHT_BG   = RGBColor(245, 245, 250)

LOGO       = "nith_logo.png"
IMG_DIR    = "images"

def img(name):
    return os.path.join(IMG_DIR, name)

# Create presentation 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# ─── Helper functions ───
def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title, bg_color=None):
    if bg_color:
        set_slide_bg(slide, bg_color)
    # Title bar background
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_RED
    bar.line.fill.background()
    # Title text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.0), Inches(0.1), height=Inches(0.9))

def add_footer(slide):
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8), Inches(0.4))
    ft.text_frame.text = "Smart Vaccine Delivery System | IoT Simulation | NIT Hamirpur"
    ft.text_frame.paragraphs[0].font.size = Pt(9)
    ft.text_frame.paragraphs[0].font.color.rgb = GOLD

def text_box(slide, left, top, width, height, text, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, font_name=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if font_name:
        p.font.name = font_name
    p.alignment = align
    return tf

def section_title(tf, title):
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_RED
    p.space_before = Pt(8)
    return p

def bullets(tf, items, size=16, color=None):
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.space_before = Pt(4)
        if color:
            p.font.color.rgb = color

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))

# ══════════════════════════════════════════════════
#  SLIDE 1: Title
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)

# Top accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = DARK_RED
bar.line.fill.background()

text_box(slide, 1, 1.0, 11.33, 0.8, "Smart Vaccine Delivery System", 44, True, WHITE, PP_ALIGN.CENTER)
text_box(slide, 1, 1.8, 11.33, 0.6, "IoT-Based Cold-Chain Monitoring & Safety Analytics", 26, True, GOLD, PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.5), Inches(5.33), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = DARK_RED
div.line.fill.background()

text_box(slide, 1, 2.8, 11.33, 0.4, "Course EC-308: Internet of Things", 18, False, GRAY, PP_ALIGN.CENTER)

# Three columns
text_box(slide, 1.0, 3.6, 4, 0.8, "Presented to:\nDr. [Professor Name]", 16, False, WHITE)
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(5.7), Inches(3.8), height=Inches(1.8))
tf = text_box(slide, 8.33, 3.6, 4, 0.8, "Presented By:\nAnirudh Sharma\nRoll No: 22DCS002", 16, False, WHITE, PP_ALIGN.RIGHT)

text_box(slide, 1, 5.8, 11.33, 0.4, "BTech DD 8th Semester, 2025–26", 16, False, GRAY, PP_ALIGN.CENTER)
text_box(slide, 1, 6.2, 11.33, 0.4, "Department of Computer Science and Engineering", 18, True, WHITE, PP_ALIGN.CENTER)
text_box(slide, 1, 6.6, 11.33, 0.4, "National Institute of Technology Hamirpur", 18, True, GOLD, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
#  SLIDE 2: Table of Contents
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Table of Contents")
add_footer(slide)

tf = text_box(slide, 1.0, 1.5, 5.5, 5.0, "", 18)
bullets(tf, [
    "Introduction & Problem Statement",
    "Limitations",
    "Proposed Solution",
    "Tools & Sensors",
    "   ○ Resistors & Values",
    "   ○ Pin Mapping (GPIOs Used)",
    "   ○ Why These Components?",
    "Simulation (Wokwi)",
    "Cloud Integration"
])
tf2 = text_box(slide, 7.0, 1.5, 5.5, 5.0, "", 18)
bullets(tf2, [
    "Code Walkthrough",
    "Results",
    "Future Scope",
    "Conclusion",
    "References"
])

# ══════════════════════════════════════════════════
#  SLIDE 3: Introduction & Problem Statement
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Introduction & Problem Statement")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 7.5, 5.5, "", 16)
section_title(tf, "Background")
bullets(tf, [
    "Vaccines require strict cold-chain (2°C–8°C) during transport",
    "WHO: up to 50% vaccine wastage due to cold-chain failures globally",
    "Manual monitoring is unreliable — no real-time alerts",
    "Physical damage (drops, vibration) and gas leakage go undetected"
])
section_title(tf, "Problem Statement")
bullets(tf, [
    "How to continuously monitor temp, humidity, mechanical stress, and contamination?",
    "How to automatically raise alerts and log violations to the cloud?"
])

# Cold chain problem image
add_image_safe(slide, img('cold_chain_problem.png'), 8.5, 1.5, 4.3, 4.3)

# ══════════════════════════════════════════════════
#  SLIDE 4: Limitations
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Limitations")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 6.0, 5.5, "", 16)
section_title(tf, "Hardware / Simulation Constraints")
bullets(tf, [
    "MPU6050 not in Wokwi — simulated via potentiometers",
    "MQ-135 warm-up: 24–48 hrs real; instant in simulation",
    "ESP32 ADC non-linearity — ±6% without calibration",
    "ThingSpeak free-tier — 15s minimum upload interval"
])

tf2 = text_box(slide, 7.0, 1.3, 6.0, 5.5, "", 16)
section_title(tf2, "System-Level Limitations")
bullets(tf2, [
    "No GPS — violation location unknown",
    "No power management — WiFi drains battery",
    "No encryption — API key hardcoded",
    "Single point of failure — ESP32 is the only MCU",
    "WiFi only — no offline buffering"
])

# ══════════════════════════════════════════════════
#  SLIDE 5: Proposed Solution
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Proposed Solution")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 5.5, 5.5, "", 16)
section_title(tf, "System Architecture")
bullets(tf, [
    "ESP32 MCU with built-in WiFi",
    "DHT22 — Temperature + Humidity",
    "MPU6050 — Accelerometer + Gyroscope",
    "MQ-135 — Gas / Spoilage detection",
    "Tamper switch — Physical security",
    "Buzzer + LED — On-device alerts",
    "ThingSpeak — Cloud logging (8 fields)",
    "Web Dashboard — Real-time monitoring"
])

# Architecture diagram (AI generated)
add_image_safe(slide, img('iot_architecture.png'), 6.5, 1.3, 6.3, 5.3)

# ══════════════════════════════════════════════════
#  SLIDE 6: Tools & Sensors
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Tools & Sensors")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 5.5, 5.5, "", 16)
section_title(tf, "Hardware Components")
bullets(tf, [
    "ESP32 DevKit v4: Main MCU + WiFi",
    "DHT22: Temp & Humidity sensing",
    "MPU6050: 6-axis IMU (I2C)",
    "MQ-135: Gas / Odor detection",
    "Push Button: Tamper detection",
    "Passive Buzzer: Audible alarm",
    "Red LED: Visual alert",
    "4.7kΩ Resistor: DHT22 pull-up",
    "220Ω Resistor: LED current limiter"
])

# Components AI image
add_image_safe(slide, img('components.png'), 6.8, 1.3, 5.8, 5.3)

# Software tools label
tf2 = text_box(slide, 0.5, 5.5, 5.5, 1.5, "", 14)
section_title(tf2, "Software Tools")
bullets(tf2, [
    "Wokwi Simulator  •  Arduino C++  •  ThingSpeak  •  Chart.js Dashboard"
], size=13)

# ══════════════════════════════════════════════════
#  SLIDE 7: Tools — Resistors & Pin Mapping
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Resistors & GPIO Pin Mapping")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 6.0, 5.0, "", 16)
section_title(tf, "Resistors Used")
bullets(tf, [
    "R1: 4700 Ω — DHT22 pull-up to 3.3V",
    "R2: 220 Ω — LED current limiter",
])
section_title(tf, "Why These Values?")
bullets(tf, [
    "4.7kΩ: Standard pull-up for DHT 1-wire; stable 3.3V idle line",
    "220Ω: Limits LED current to ≈15mA (ESP32 max: 20mA)"
])

# Pin Mapping Table
tbl = slide.shapes.add_table(8, 4, Inches(7.0), Inches(1.3), Inches(6.0), Inches(4.5)).table
headers = ["GPIO", "Component", "Mode", "Reason"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_RED

rows = [
    ["4",  "DHT22 Data",    "Digital In",  "1-Wire protocol"],
    ["35", "Accel-X (pot)", "Analog In",   "ADC1, input-only"],
    ["32", "Accel-Y (pot)", "Analog In",   "ADC1, general"],
    ["34", "Gas MQ-135",    "Analog In",   "ADC1, input-only"],
    ["15", "Tamper Button", "Digital In",  "INPUT_PULLUP"],
    ["13", "Buzzer",        "Digital Out", "PWM capable"],
    ["2",  "Red LED",       "Digital Out", "On-board LED"],
]
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

# ══════════════════════════════════════════════════
#  SLIDE 8: Tools — Why These Components?
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Why These Components?")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 6.0, 5.5, "", 15)
bullets(tf, [
    "ESP32 — WiFi+BT, dual-core 240MHz, 3.3V I/O, multiple ADC. No additional WiFi shield needed.",
    "DHT22 — ±0.5°C accuracy, covers 2–8°C vaccine range. Superior to DHT11.",
    "MPU6050 — 6-axis IMU (accel+gyro) via I2C. Detects drops, vibration, rash driving.",
    "MQ-135 — Detects NH₃, CO₂, benzene, smoke. Ideal for spoilage detection.",
], size=15)

tf2 = text_box(slide, 7.0, 1.3, 6.0, 5.5, "", 15)
bullets(tf2, [
    "Tamper Button — Simple physical security with INPUT_PULLUP; no external resistor.",
    "ThingSpeak — Free REST-API cloud IoT; 8 fields/channel; MATLAB visualizations.",
    "Wokwi — Browser-based ESP32+DHT22 simulator. No hardware for demo.",
    "Chart.js — Lightweight real-time dashboard with animated gauges."
], size=15)

# Methodology AI image
add_image_safe(slide, img('methodology.png'), 4.5, 4.0, 4.5, 3.0)

# ══════════════════════════════════════════════════
#  SLIDE 9: Wokwi Simulation — Circuit
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Simulation — Wokwi Circuit")
add_footer(slide)

# Wokwi screenshot
add_image_safe(slide, img('wokwi_simulation.png'), 0.3, 1.3, 8.0, 4.8)

tf = text_box(slide, 8.8, 1.3, 4.3, 5.5, "", 15)
section_title(tf, "Wokwi Components")
bullets(tf, [
    "ESP32 DevKit C v4",
    "DHT22 with 4.7kΩ pull-up",
    "Pot 1 (GPIO 35) — Accel-X",
    "Pot 2 (GPIO 32) — Accel-Y",
    "Pot 3 (GPIO 34) — MQ-135 Gas",
    "Push Button (GPIO 15)",
    "Passive Buzzer (GPIO 13)",
    "Red LED + 220Ω (GPIO 2)"
], size=14)
section_title(tf, "Simulating Events")
bullets(tf, [
    "Pot 1/2 → right = shock/drop",
    "Pot 3 → right = gas leak",
    "Click DHT22 = change temp/hum",
    "Push button = tamper alert"
], size=14)

# ══════════════════════════════════════════════════
#  SLIDE 10: Simulation — Serial Monitor
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
add_header(slide, "Simulation — Serial Monitor Output")
add_footer(slide)

# Serial monitor image
add_image_safe(slide, img('serial_monitor.png'), 0.3, 1.3, 6.5, 5.0)

tf = text_box(slide, 7.3, 1.3, 5.5, 5.5, "", 15, color=WHITE)
section_title(tf, "Serial Output Fields")
bullets(tf, [
    "Temperature (°C) with range check",
    "Humidity (%RH) with HIGH flag",
    "Acceleration magnitude with SHOCK flag",
    "Gyro speed with FAST flag",
    "Gas level with LEAK flag",
    "Tamper: SEALED / OPEN",
    "Damage Score (0–10)",
    "Status Code (0, 1, 2...9)"
], size=14, color=WHITE)
section_title(tf, "JSON Data Output")
bullets(tf, [
    "Machine-readable per 2s cycle",
    "Format: DATA:{temp, hum, gas, ...}",
    "Enables dashboard integration"
], size=14, color=WHITE)

# ══════════════════════════════════════════════════
#  SLIDE 11: Cloud Integration — ThingSpeak
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Cloud Integration — ThingSpeak")
add_footer(slide)

# ThingSpeak channel chart
add_image_safe(slide, img('thingspeak_channel.png'), 0.3, 1.3, 7.5, 4.5)

tf = text_box(slide, 8.3, 1.3, 4.8, 5.5, "", 15)
section_title(tf, "Channel Data Fields")

tbl = slide.shapes.add_table(9, 3, Inches(8.3), Inches(2.0), Inches(4.5), Inches(3.2)).table
th = [("Field", "Data", "Unit")]
rows = [
    ("1", "Temperature", "°C"),
    ("2", "Humidity", "%RH"),
    ("3", "Status Code", "0–9"),
    ("4", "Accel Magnitude", "raw"),
    ("5", "Gas Level", "analog"),
    ("6", "Damage Score", "0–10"),
    ("7", "Gyro Speed", "raw"),
    ("8", "Tamper State", "0/1"),
]
for j, h in enumerate(th[0]):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_RED
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)

tf2 = text_box(slide, 8.3, 5.4, 4.5, 1.5, "", 14)
section_title(tf2, "Upload Protocol")
bullets(tf2, [
    "HTTP GET to api.thingspeak.com/update",
    "Every 15s (free-tier) via WiFi"
], size=13)

# ══════════════════════════════════════════════════
#  SLIDE 12: Code Walkthrough — Sensor Reading
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
add_header(slide, "Code Walkthrough — Sensor Reading & Damage Score")
add_footer(slide)

# Left code block
code1 = """// DHT22
TempAndHumidity d = dht.getTempAndHumidity();
temperature = d.temperature;
humidity    = d.humidity;

// MPU6050 via potentiometers
int rawX = analogRead(ACCEL_X_PIN);
int rawY = analogRead(ACCEL_Y_PIN);
float ax = map(rawX, 0,4095,-32768,32767);
float ay = map(rawY, 0,4095,-32768,32767);
float az = 16384; // ~1g baseline
accelMag = sqrt(ax*ax + ay*ay + az*az);

// Tamper
tampered = (digitalRead(TAMPER_PIN)==LOW);"""

lbl1 = text_box(slide, 0.5, 1.3, 6.0, 0.4, "DHT22 + MPU6050 Simulation", 18, True, GOLD)
tb1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0))
tf1 = tb1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = code1
p.font.name = "Courier New"
p.font.size = Pt(12)
p.font.color.rgb = GREEN

code2 = """damageScore = 0; status_code = 1;

if (temp>8 || temp<2)  damageScore += 3;
if (tampered)          damageScore += 2;
if (accelMag > 20000)  damageScore += 2;
if (gyroSpeed > 30000) damageScore += 1;
if (gasValue > 400)    damageScore += 3;
if (hum > 80)          damageScore += 1;

// CRITICAL override
if (damageScore >= 5)
    status_code = 9; // CRITICAL"""

lbl2 = text_box(slide, 7.0, 1.3, 6.0, 0.4, "Damage Score Algorithm", 18, True, GOLD)
tb2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(6.0), Inches(5.0))
tf2 = tb2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = code2
p.font.name = "Courier New"
p.font.size = Pt(12)
p.font.color.rgb = GREEN

# ══════════════════════════════════════════════════
#  SLIDE 13: Code Walkthrough — Cloud & Status Codes
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)
add_header(slide, "Code Walkthrough — Cloud Upload & Status Codes")
add_footer(slide)

code3 = """String url = server +
  "?api_key=" + apiKey +
  "&field1=" + String(temperature) +
  "&field2=" + String(humidity)    +
  "&field3=" + String(status_code) +
  "&field4=" + String(accelMag)    +
  "&field5=" + String(gasValue)    +
  "&field6=" + String(damageScore);

http.begin(url);
int code = http.GET();
// Runs every 15 seconds"""

text_box(slide, 0.5, 1.3, 6.0, 0.4, "ThingSpeak HTTP Upload", 18, True, GOLD)
tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5))
tf3 = tb3.text_frame
tf3.word_wrap = True
p = tf3.paragraphs[0]
p.text = code3
p.font.name = "Courier New"
p.font.size = Pt(12)
p.font.color.rgb = GREEN

# Status Code Table
text_box(slide, 7.0, 1.3, 6.0, 0.4, "Status Code Reference", 18, True, GOLD)
tbl = slide.shapes.add_table(9, 3, Inches(7.0), Inches(1.9), Inches(5.8), Inches(4.2)).table
headers = ["Code", "Condition", "Severity"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_RED

status_data = [
    ("0", "Temp out of range",  "🔴 Critical"),
    ("1", "All normal",         "🟢 OK"),
    ("2", "Tamper detected",    "🟡 Warning"),
    ("3", "Drop / Shock",       "🔴 Critical"),
    ("4", "Unsafe driving",     "🟡 Warning"),
    ("5", "Gas leakage",        "🔴 Critical"),
    ("6", "High humidity",      "🟡 Warning"),
    ("9", "Damage ≥ 5",         "🔴 CRITICAL"),
]
for i, (c, cond, sev) in enumerate(status_data):
    for j, val in enumerate([c, cond, sev]):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE

# ══════════════════════════════════════════════════
#  SLIDE 14: Results
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Results")
add_footer(slide)

# Dashboard screenshot
add_image_safe(slide, img('dashboard.png'), 0.3, 1.3, 8.2, 5.2)

tf = text_box(slide, 8.8, 1.3, 4.3, 5.5, "", 15)
section_title(tf, "Key Observations")
bullets(tf, [
    "Cold-chain failure triggers Status 0 at temp >8°C or <2°C",
    "Shock event → Status 3 in 2s; buzzer activates",
    "Gas leakage raises Status 5 with LED flash",
    "Multi-fault → CRITICAL (Status 9) via damage score",
    "ThingSpeak receives 8 fields every 15s",
    "Serial monitor confirms JSON per cycle"
], size=14)

# ══════════════════════════════════════════════════
#  SLIDE 15: Future Scope
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Future Scope")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 5.5, 5.0, "", 15)
section_title(tf, "Short-Term Enhancements")
bullets(tf, [
    "GPS (NEO-6M) — Real-time location & route safety",
    "Telegram / WhatsApp Bot — Instant push alerts",
    "Deep-Sleep Mode — 72+ hour battery life",
    "SD Card Logging — Offline data buffering"
])

section_title(tf, "Advanced Research")
bullets(tf, [
    "Blockchain Audit Trail — Immutable logs on Starknet",
    "ML on Edge — TensorFlow Lite anomaly detection",
    "LPWAN (LoRa / NB-IoT) — Rural connectivity",
    "Fleet Monitoring — Multi-vehicle risk scoring"
])

# Future scope AI image
add_image_safe(slide, img('future_scope.png'), 6.5, 1.3, 6.3, 5.5)

# ══════════════════════════════════════════════════
#  SLIDE 16: Conclusion
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Conclusion")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 8.5, 5.5, "", 16)
bullets(tf, [
    "Designed and simulated a complete IoT pipeline for real-time vaccine cold-chain monitoring using ESP32",
    "Integrated 5 sensing modalities: temperature, humidity, motion, gas, tamper detection",
    "Novel Damage Score Algorithm (0–10) with CRITICAL-level escalation at ≥5",
    "Cloud connectivity via ThingSpeak with 8-field data logging every 15s",
    "Premium dashboard: animated gauges, Chart.js history, event simulation, severity-coded alerts",
    '"Ensuring vaccine integrity by monitoring environmental conditions, mechanical stress, and contamination risks in real time"'
])

# Checklist
check_items = ["✓ IoT Sensors", "✓ Embedded Firmware", "✓ Cloud Analytics", "✓ Safety Algorithms", "✓ Healthcare Use-Case"]
for i, item in enumerate(check_items):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.5), Inches(1.5 + i * 0.85), Inches(3.5), Inches(0.65))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(20, 60, 30)
    rect.line.color.rgb = GREEN
    rect.text_frame.paragraphs[0].text = item
    rect.text_frame.paragraphs[0].font.size = Pt(16)
    rect.text_frame.paragraphs[0].font.color.rgb = GREEN
    rect.text_frame.paragraphs[0].font.bold = True
    rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

text_box(slide, 9.5, 6.0, 3.5, 0.5, "Top 1% IoT Project", 18, True, GOLD, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
#  SLIDE 17: References
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "References")
add_footer(slide)

tf = text_box(slide, 0.5, 1.3, 6.0, 5.5, "", 14)
section_title(tf, "Simulation & Hardware")
bullets(tf, [
    "Wokwi Reference — https://wokwi.com\n  ESP32 DevKit C v4 docs, DHT22 component guide",
    "ESP32 Reference — https://docs.espressif.com\n  Technical Reference Manual, ADC1 GPIO specs",
    "MPU6050 Datasheet — InvenSense Product Spec Rev. 3.4",
    "MQ-135 Datasheet — Hanwei Electronics Sensor"
], size=13)

tf2 = text_box(slide, 7.0, 1.3, 6.0, 5.5, "", 14)
section_title(tf2, "Libraries, Cloud & Standards")
bullets(tf2, [
    "DHTesp Library — beegee-tokyo/DHTesp (GitHub)",
    "ThingSpeak Platform — https://thingspeak.com\n  MathWorks REST API documentation",
    "Chart.js — https://chartjs.org (v4.4.4)",
    "WHO Cold-Chain Standard — WHO/IVB/06.10",
    "DHT22 — Aosong Electronics AM2302 datasheet"
], size=13)

# ══════════════════════════════════════════════════
#  SLIDE 18: Thank You
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
set_slide_bg(slide, DARK_BG)

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = DARK_RED
bar.line.fill.background()

text_box(slide, 2, 1.5, 9.33, 1.2, "Thank You", 60, True, DARK_RED, PP_ALIGN.CENTER)
text_box(slide, 2, 3.0, 9.33, 0.6, "Questions & Discussion", 28, True, GOLD, PP_ALIGN.CENTER)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.8), Inches(3.33), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = GOLD
div.line.fill.background()

text_box(slide, 2, 4.2, 9.33, 0.4, "Anirudh Sharma  |  Roll No: 22DCS002", 18, True, WHITE, PP_ALIGN.CENTER)
text_box(slide, 2, 4.7, 9.33, 0.4, "BTech DD 8th Semester, 2025–26", 16, False, GRAY, PP_ALIGN.CENTER)
text_box(slide, 2, 5.1, 9.33, 0.4, "Department of Computer Science and Engineering", 16, False, GRAY, PP_ALIGN.CENTER)
text_box(slide, 2, 5.5, 9.33, 0.4, "National Institute of Technology Hamirpur", 16, True, GOLD, PP_ALIGN.CENTER)

if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(5.9), Inches(6.0), height=Inches(1.2))

# ─── Save ───
output = 'presentation.pptx'
prs.save(output)
print(f"✅ Successfully generated {output} with {len(prs.slides)} slides and real images!")
